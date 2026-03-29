#!/usr/bin/env python3
"""
Ingest PDF, Word, and PowerPoint files into the project's knowledge-base CSV format.
"""

from __future__ import annotations

import argparse
import csv
import fnmatch
import hashlib
import json
import re
import shutil
import subprocess
import unicodedata
import zipfile
from datetime import datetime, timezone
from pathlib import Path

try:
    from docx import Document as DocxDocument
    from pypdf import PdfReader
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover
    raise SystemExit(
        "Missing dependencies. Run: pip install -r scripts/requirements-document-ingest.txt"
    ) from exc

try:  # Optional higher-quality PDF extractor
    import fitz
except ImportError:  # pragma: no cover
    fitz = None

try:  # Optional higher-quality PDF extractor
    import pdfplumber
except ImportError:  # pragma: no cover
    pdfplumber = None

try:  # Optional OCR fallback
    from rapidocr_onnxruntime import RapidOCR
except ImportError:  # pragma: no cover
    RapidOCR = None


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".doc", ".ppt"}
DIRECT_EXTENSIONS = {".pdf", ".docx", ".pptx"}
ZIP_EXTENSION = ".zip"
CONVERSION_TARGETS = {".doc": ".docx", ".ppt": ".pptx"}

KB_FIELDNAMES = [
    "id",
    "title",
    "category",
    "subcategory",
    "lab_type",
    "risk_level",
    "hazard_types",
    "scenario",
    "question",
    "answer",
    "steps",
    "ppe",
    "forbidden",
    "disposal",
    "first_aid",
    "emergency",
    "legal_notes",
    "references",
    "source_type",
    "source_title",
    "source_org",
    "source_version",
    "source_date",
    "source_url",
    "last_updated",
    "reviewer",
    "status",
    "tags",
    "language",
]

NOISE_PATTERNS = [
    r"^\d+$",
    r"^第?\s*\d+\s*页$",
    r"^Page\s+\d+$",
    r"^版权所有",
    r"^Copyright",
    r"^目录$",
    r"^Contents$",
]

LAB_KEYWORDS = {
    "化学": "化学",
    "危化": "化学",
    "试剂": "化学",
    "酸": "化学",
    "碱": "化学",
    "电气": "电气",
    "高压": "电气",
    "触电": "电气",
    "生物": "生物",
    "菌": "生物",
    "消防": "通用",
    "火灾": "通用",
    "实验室": "通用",
}

HAZARD_KEYWORDS = {
    "高压": "高压",
    "易燃": "易燃",
    "腐蚀": "腐蚀性",
    "辐射": "辐射",
    "X射线": "辐射",
    "气瓶": "高压气体",
    "火灾": "火灾",
    "触电": "触电",
}

PDF_PYPDF_SHORT_LINE_RATIO_THRESHOLD = 0.18
PDF_PYPDF_SINGLE_CHAR_RATIO_THRESHOLD = 0.08
PDF_OCR_SHORT_LINE_RATIO_THRESHOLD = 0.12
PDF_OCR_SINGLE_CHAR_RATIO_THRESHOLD = 0.05
PDF_OCR_MIN_CHAR_COUNT = 1000
PDF_OCR_RENDER_SCALE = 2.0
PDF_RARE_HAN_RANGE_START = 0x7280
PDF_RARE_HAN_RANGE_END = 0x72FF
PDF_QUALITY_KEYWORDS = [
    "范围",
    "总则",
    "要求",
    "规则",
    "规范",
    "标准",
    "安全",
    "实验室",
    "附录",
]

_rapidocr_engine: object | None = None


def now_iso() -> str:
    return datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds")


def clean_line(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def normalize_text(text: str) -> str:
    text = text.replace("\r", "\n")
    text = text.replace("\x00", "")
    raw_lines = [clean_line(line) for line in text.splitlines()]
    lines: list[str] = []
    previous = ""
    for line in raw_lines:
        if not line:
            continue
        if any(re.search(pattern, line, re.IGNORECASE) for pattern in NOISE_PATTERNS):
            continue
        if line == previous:
            continue
        previous = line
        lines.append(line)
    return "\n".join(lines).strip()


def normalize_pdf_line(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\u3000", " ")
    text = text.replace("\ufeff", "")
    text = text.replace("\u00ad", "")
    text = "".join(ch for ch in text if unicodedata.category(ch) != "Co")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def normalize_pdf_compact(text: str) -> str:
    return re.sub(r"\s+", "", unicodedata.normalize("NFKC", text or ""))


def normalize_pdf_page_head(lines: list[str], limit: int = 8) -> str:
    return normalize_pdf_compact("".join(lines[:limit]))


def looks_like_page_number(line: str) -> bool:
    compact = line.strip()
    if re.fullmatch(r"[0-9]{1,4}", compact):
        return True
    if re.fullmatch(r"[0-9]{1,4}\s*(?:/|of)\s*[0-9]{1,4}", compact, re.IGNORECASE):
        return True
    if re.fullmatch(r"[IVXLCM]{1,8}", compact):
        return True
    if re.fullmatch(r"第\s*\d+\s*页", compact):
        return True
    return False


def looks_like_standard_code(line: str) -> bool:
    compact = normalize_pdf_compact(line)
    if re.search(r"\b(?:GB|GB/T|IEC|ISO|AQ|JJG|JJF|DIN|JB)[A-Z0-9/\-.—]{2,}\d", compact, re.IGNORECASE):
        return True
    return bool(re.search(r"(?:^|\D)\d{4,6}[—-]\d{4}(?:\D|$)", compact))


def looks_like_toc_line(line: str) -> bool:
    if "…" in line or "..." in line or "···" in line:
        return True
    return bool(re.search(r"^.{4,}(?:\s+|\.{2,}|…{2,}|···)[IVXLCM\d]{1,6}$", line.strip()))


def looks_like_front_matter_heading(line: str) -> bool:
    compact = normalize_pdf_compact(line).lower()
    return compact in {"前言", "引言", "序言", "目录", "目次", "contents"}


def page_has_front_matter_heading(lines: list[str]) -> bool:
    head = normalize_pdf_page_head(lines, limit=10).lower()
    return any(token in head[:32] for token in ("前言", "引言", "序言", "目录", "目次", "contents"))


def looks_like_section_heading(line: str) -> bool:
    compact = line.strip()
    patterns = [
        r"^\d+(\.\d+)*\s+\S+",
        r"^第[一二三四五六七八九十百]+[章节部分]\s*.*",
        r"^附录[A-Z一二三四五六七八九十]?\s*.*",
        r"^(图|表)\s*\d+.*",
        r"^前\s*言$",
    ]
    return any(re.match(pattern, compact) for pattern in patterns)


def is_pdf_cover_page(page_index: int, lines: list[str]) -> bool:
    if page_index > 2 or len(lines) < 5:
        return False
    joined = "\n".join(lines[:40])
    compact = normalize_pdf_page_head(lines, limit=40)
    compact_joined = normalize_pdf_compact(joined)
    cover_signals = [
        "中华人民共和国国家标准",
        "国家标准",
        "发布",
        "实施",
    ]
    hits = sum(
        1
        for signal in cover_signals
        if normalize_pdf_compact(signal) in compact or normalize_pdf_compact(signal) in compact_joined
    )
    has_standard_marker = any(looks_like_standard_code(line) for line in lines[:20]) or bool(
        re.search(r"\d{4,6}[—-]\d{4}", compact)
    )
    short_line_count = sum(1 for line in lines[:20] if len(normalize_pdf_compact(line)) <= 8)
    has_short_line_layout = short_line_count >= 6 or (page_index == 1 and short_line_count >= 4)
    has_publish_schedule = bool(re.search(r"20\d{6}发布.*20\d{6}实施", compact_joined))
    return hits >= 3 and has_standard_marker and (has_short_line_layout or has_publish_schedule)


def is_pdf_toc_page(lines: list[str]) -> bool:
    if not lines:
        return False
    joined = "\n".join(lines[:30])
    compact = normalize_pdf_compact(joined).lower()
    if any(token in compact for token in ("目录", "目次", "contents")):
        return True
    toc_line_count = sum(1 for line in lines[:80] if looks_like_toc_line(line))
    return toc_line_count >= 8


def find_pdf_body_start_page(page_lines: list[list[str]]) -> int | None:
    search_limit = min(20, len(page_lines))
    for index in range(search_limit):
        lines = page_lines[index]
        if is_pdf_cover_page(index + 1, lines) or is_pdf_toc_page(lines):
            continue
        compact = normalize_pdf_compact("\n".join(lines[:80]))
        if re.search(r"(^|[^0-9])1范围", compact):
            return index + 1
        if re.search(r"(^|[^0-9])1总则", compact):
            return index + 1
        if re.search(r"(^|[^0-9])1scope", compact, re.IGNORECASE):
            return index + 1
        if page_has_front_matter_heading(lines):
            continue
        paired_headings = (
            normalize_pdf_compact(f"{left}{right}")
            for left, right in zip(lines[:20], lines[1:21])
        )
        if any(re.fullmatch(r"1(?:范围|总则|scope)", pair, re.IGNORECASE) for pair in paired_headings):
            return index + 1
    return None


def collect_repeated_pdf_lines(page_lines: list[list[str]]) -> set[str]:
    counts: dict[str, int] = {}
    for lines in page_lines:
        unique_lines = set(
            line
            for line in lines
            if line
            and len(line) <= 40
            and (looks_like_page_number(line) or looks_like_standard_code(line) or len(line) <= 20)
        )
        for line in unique_lines:
            counts[line] = counts.get(line, 0) + 1

    min_repeat = max(3, len(page_lines) // 5)
    return {line for line, count in counts.items() if count >= min_repeat}


def merge_pdf_wrapped_lines(lines: list[str]) -> list[str]:
    merged: list[str] = []
    for line in lines:
        if not merged:
            merged.append(line)
            continue

        previous = merged[-1]
        if (
            len(previous) < 80
            and len(line) < 80
            and not looks_like_section_heading(previous)
            and not looks_like_section_heading(line)
            and not re.search(r"[。！？；：:】）》\]\)]$", previous)
            and not re.match(r"^(\d+(\.\d+)*|第[一二三四五六七八九十]+[章节部分]|附录|注[:：]|表\d+|图\d+|[-—])", line)
        ):
            merged[-1] = f"{previous}{line}"
        else:
            merged.append(line)
    return merged


def get_rapidocr_engine() -> object:
    global _rapidocr_engine
    if RapidOCR is None:
        raise RuntimeError("RapidOCR is not installed")
    if _rapidocr_engine is None:
        _rapidocr_engine = RapidOCR()
    return _rapidocr_engine


def normalize_pdf_lines_from_text(text: str) -> list[str]:
    return [line for line in (normalize_pdf_line(item) for item in text.splitlines()) if line]


def extract_pdf_pages_pypdf(path: Path) -> list[list[str]]:
    reader = PdfReader(str(path))
    return [normalize_pdf_lines_from_text(page.extract_text() or "") for page in reader.pages]


def extract_pdf_pages_pymupdf(path: Path) -> list[list[str]]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is not installed")
    with fitz.open(str(path)) as document:
        return [normalize_pdf_lines_from_text(page.get_text("text")) for page in document]


def extract_pdf_pages_pdfplumber(path: Path) -> list[list[str]]:
    if pdfplumber is None:
        raise RuntimeError("pdfplumber is not installed")
    with pdfplumber.open(str(path)) as document:
        return [normalize_pdf_lines_from_text(page.extract_text() or "") for page in document.pages]


def extract_pdf_pages_ocr(path: Path, render_scale: float = PDF_OCR_RENDER_SCALE) -> list[list[str]]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for OCR rendering")
    engine = get_rapidocr_engine()
    with fitz.open(str(path)) as document:
        pages: list[list[str]] = []
        for page in document:
            pixmap = page.get_pixmap(matrix=fitz.Matrix(render_scale, render_scale), alpha=False)
            result, _ = engine(pixmap.tobytes("png"))
            raw_lines = []
            for item in result or []:
                if isinstance(item, (list, tuple)) and len(item) >= 2:
                    raw_lines.append(str(item[1]))
            pages.append(normalize_pdf_lines_from_text("\n".join(raw_lines)))
    return pages


def count_pdf_prefix_rare_han(text: str, limit: int = 200) -> int:
    prefix = normalize_text(text)[:limit]
    return sum(1 for ch in prefix if PDF_RARE_HAN_RANGE_START <= ord(ch) <= PDF_RARE_HAN_RANGE_END)


def count_pdf_keyword_hits(text: str, limit: int = 1200) -> int:
    prefix = normalize_text(text)[:limit]
    return sum(prefix.count(keyword) for keyword in PDF_QUALITY_KEYWORDS)


def clean_pdf_pages(raw_pages: list[list[str]], special_rule: dict | None = None) -> tuple[str, dict]:
    repeated_lines = collect_repeated_pdf_lines(raw_pages)
    body_start_page = find_pdf_body_start_page(raw_pages)
    if special_rule and special_rule.get("body_start_page"):
        body_start_page = min(special_rule["body_start_page"], len(raw_pages))
    toc_scan_limit = body_start_page - 1 if body_start_page else min(8, len(raw_pages))
    special_skip_pages = set(special_rule.get("skip_pages", [])) if special_rule else set()
    kept_pages: list[str] = []
    skipped_pages: list[dict] = []

    for index, lines in enumerate(raw_pages, start=1):
        if not lines:
            skipped_pages.append({"page": index, "reason": "empty_page"})
            continue
        if is_pdf_cover_page(index, lines):
            skipped_pages.append({"page": index, "reason": "cover_page"})
            continue
        if index <= toc_scan_limit and is_pdf_toc_page(lines):
            skipped_pages.append({"page": index, "reason": "toc_page"})
            continue
        if body_start_page and index < body_start_page:
            skipped_pages.append({"page": index, "reason": "front_matter_page"})
            continue
        if index in special_skip_pages:
            skipped_pages.append({"page": index, "reason": "special_rule_skip"})
            continue

        filtered = [
            line
            for line in lines
            if line not in repeated_lines and not looks_like_page_number(line)
        ]
        filtered = [line for line in filtered if len(line) > 1]
        filtered = merge_pdf_wrapped_lines(filtered)
        page_text = normalize_text("\n".join(filtered))

        if len(page_text) < 40:
            skipped_pages.append({"page": index, "reason": "too_short_after_cleaning"})
            continue

        kept_pages.append(page_text)

    payload = {
        "page_count": len(raw_pages),
        "kept_page_count": len(kept_pages),
        "body_start_page": body_start_page,
        "skipped_pages": skipped_pages,
        "repeated_line_count": len(repeated_lines),
    }
    if special_rule:
        payload["pdf_special_rule"] = {
            "rule_id": special_rule["rule_id"],
            "path_pattern": special_rule["path_pattern"],
            "force_ocr": special_rule["force_ocr"],
            "body_start_page": special_rule.get("body_start_page"),
            "skip_pages": special_rule.get("skip_pages", []),
            "notes": special_rule.get("notes", ""),
        }
    return "\n\n".join(kept_pages), payload


def summarize_pdf_candidate(method: str, raw_pages: list[list[str]], cleaned_text: str, clean_meta: dict) -> dict:
    raw_lines = [line for page in raw_pages for line in page if line]
    raw_line_count = len(raw_lines) or 1
    short_line_count = sum(1 for line in raw_lines if len(line) <= 2)
    single_char_count = sum(1 for line in raw_lines if len(line) == 1)
    short_line_ratio = short_line_count / raw_line_count
    single_char_ratio = single_char_count / raw_line_count
    char_count = len(cleaned_text)
    keyword_hits = count_pdf_keyword_hits(cleaned_text)
    prefix_rare_han_count = count_pdf_prefix_rare_han(cleaned_text)

    score = float(min(char_count, 40000))
    score -= short_line_ratio * 25000
    score -= single_char_ratio * 15000
    score -= prefix_rare_han_count * 180
    score += min(keyword_hits, 20) * 60
    if clean_meta.get("body_start_page"):
        score += 600
    if clean_meta.get("kept_page_count", 0) <= 1:
        score -= 500

    return {
        "method": method,
        "score": round(score, 2),
        "raw_page_count": len(raw_pages),
        "raw_line_count": raw_line_count,
        "char_count": char_count,
        "short_line_ratio": round(short_line_ratio, 4),
        "single_char_ratio": round(single_char_ratio, 4),
        "keyword_hits": keyword_hits,
        "prefix_rare_han_count": prefix_rare_han_count,
        "body_start_page": clean_meta.get("body_start_page"),
        "kept_page_count": clean_meta.get("kept_page_count", 0),
    }


def build_pdf_candidate(method: str, raw_pages: list[list[str]], special_rule: dict | None = None) -> dict:
    cleaned_text, clean_meta = clean_pdf_pages(raw_pages, special_rule=special_rule)
    summary = summarize_pdf_candidate(method, raw_pages, cleaned_text, clean_meta)
    return {
        "method": method,
        "raw_pages": raw_pages,
        "text": cleaned_text,
        "clean_meta": clean_meta,
        "summary": summary,
    }


def choose_best_pdf_candidate(candidates: list[dict]) -> dict:
    return max(candidates, key=lambda item: item["summary"]["score"])


def should_fallback_from_pypdf(candidate: dict) -> bool:
    summary = candidate["summary"]
    return (
        summary["char_count"] < PDF_OCR_MIN_CHAR_COUNT
        or summary["short_line_ratio"] > PDF_PYPDF_SHORT_LINE_RATIO_THRESHOLD
        or summary["single_char_ratio"] > PDF_PYPDF_SINGLE_CHAR_RATIO_THRESHOLD
    )


def should_try_ocr(best_candidate: dict, ocr_mode: str) -> bool:
    if ocr_mode == "off":
        return False
    if ocr_mode == "always":
        return True
    summary = best_candidate["summary"]
    return (
        summary["char_count"] < PDF_OCR_MIN_CHAR_COUNT
        or summary["short_line_ratio"] > PDF_OCR_SHORT_LINE_RATIO_THRESHOLD
        or summary["single_char_ratio"] > PDF_OCR_SINGLE_CHAR_RATIO_THRESHOLD
        or (summary["prefix_rare_han_count"] >= 4 and summary["char_count"] <= 8000)
    )


def patch_pdf_title_from_ocr(primary_text: str, ocr_text: str) -> tuple[str, bool]:
    primary_lines = [line for line in primary_text.splitlines() if line.strip()]
    ocr_lines = [line for line in ocr_text.splitlines() if line.strip()]
    if not primary_lines or not ocr_lines:
        return primary_text, False

    primary_title = primary_lines[0]
    ocr_title = ocr_lines[0]
    primary_rare = count_pdf_prefix_rare_han(primary_title, limit=120)
    ocr_rare = count_pdf_prefix_rare_han(ocr_title, limit=120)
    if primary_rare < 2 or ocr_rare >= primary_rare:
        return primary_text, False
    if not any(keyword in ocr_title for keyword in ("规则", "要求", "标准", "规范", "实验室", "检测")):
        return primary_text, False

    primary_lines[0] = ocr_title
    return "\n".join(primary_lines), True


def split_into_chunks(text: str, max_chars: int, overlap: int) -> list[str]:
    paragraphs = [part.strip() for part in text.split("\n") if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if not current:
            current = paragraph
            continue
        candidate = f"{current}\n{paragraph}"
        if len(candidate) <= max_chars:
            current = candidate
            continue
        chunks.append(current)
        if overlap > 0 and len(current) > overlap:
            current = current[-overlap:] + "\n" + paragraph
        else:
            current = paragraph
    if current:
        chunks.append(current)
    return [chunk.strip() for chunk in chunks if chunk.strip()]


def write_jsonl(path: Path, rows: list[dict]) -> None:
    with path.open("w", encoding="utf-8") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False) + "\n")


def write_csv(path: Path, rows: list[dict]) -> None:
    with path.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=KB_FIELDNAMES)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def merge_into_csv(path: Path, new_rows: list[dict]) -> int:
    existing_rows: list[dict] = []
    existing_ids: set[str] = set()
    if path.exists():
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.DictReader(handle)
            for row in reader:
                normalized = {field: row.get(field, "") for field in KB_FIELDNAMES}
                existing_rows.append(normalized)
                existing_ids.add(normalized["id"])

    appended = 0
    for row in new_rows:
        if row["id"] in existing_ids:
            continue
        existing_rows.append({field: row.get(field, "") for field in KB_FIELDNAMES})
        existing_ids.add(row["id"])
        appended += 1

    write_csv(path, existing_rows)
    return appended


def sha1_text(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8")).hexdigest()


def normalize_rel_path(value: str) -> str:
    return value.replace("/", "\\").strip().lower()


def parse_bool(value: str | bool | None) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return False
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def parse_positive_int(value: str | int | None) -> int | None:
    if value in (None, ""):
        return None
    try:
        parsed = int(str(value).strip())
    except ValueError:
        return None
    return parsed if parsed > 0 else None


def parse_page_spec(value: str | None) -> list[int]:
    if not value:
        return []
    pages: set[int] = set()
    for part in re.split(r"[;,，\s]+", value.strip()):
        token = part.strip()
        if not token:
            continue
        if "-" in token:
            left, right = token.split("-", 1)
            start = parse_positive_int(left)
            end = parse_positive_int(right)
            if start is None or end is None:
                continue
            if start > end:
                start, end = end, start
            pages.update(range(start, end + 1))
            continue
        single = parse_positive_int(token)
        if single is not None:
            pages.add(single)
    return sorted(pages)


def load_manifest(path: Path | None) -> dict[str, dict]:
    if not path or not path.exists():
        return {}
    result: dict[str, dict] = {}
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            key = normalize_rel_path(row.get("path") or "")
            if key:
                result[key] = row
    return result


def load_pdf_special_rules(path: Path | None) -> list[dict]:
    if not path or not path.exists():
        return []
    rules: list[dict] = []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        for index, row in enumerate(reader, start=1):
            pattern = normalize_rel_path(
                row.get("path_pattern") or row.get("path") or row.get("file_pattern") or ""
            )
            if not pattern:
                continue
            rules.append(
                {
                    "rule_id": (row.get("rule_id") or f"PDF-RULE-{index:03d}").strip(),
                    "path_pattern": pattern,
                    "force_ocr": parse_bool(row.get("force_ocr")),
                    "body_start_page": parse_positive_int(row.get("body_start_page")),
                    "skip_pages": parse_page_spec(row.get("skip_pages")),
                    "notes": (row.get("notes") or "").strip(),
                }
            )
    return rules


def match_pdf_special_rule(rel_path: str, rules: list[dict]) -> dict | None:
    normalized = normalize_rel_path(rel_path)
    file_name = Path(normalized).name
    normalized_posix = normalized.replace("\\", "/")
    for rule in rules:
        pattern = rule["path_pattern"]
        pattern_posix = pattern.replace("\\", "/")
        if (
            fnmatch.fnmatch(normalized, pattern)
            or fnmatch.fnmatch(file_name, pattern)
            or fnmatch.fnmatch(normalized_posix, pattern_posix)
        ):
            return rule
    return None


def infer_lab_type(text: str) -> str:
    for keyword, lab_type in LAB_KEYWORDS.items():
        if keyword in text:
            return lab_type
    return "通用"


def infer_hazards(text: str) -> str:
    hits: list[str] = []
    for keyword, hazard in HAZARD_KEYWORDS.items():
        if keyword in text and hazard not in hits:
            hits.append(hazard)
    return ";".join(hits) if hits else "综合安全"


def infer_subcategory(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "PDF资料"
    if suffix in {".docx", ".doc"}:
        return "Word资料"
    if suffix in {".pptx", ".ppt"}:
        return "PPT资料"
    return "文档资料"


def build_metadata(file_path: Path, input_root: Path, manifest: dict[str, dict]) -> dict:
    try:
        rel_path = str(file_path.relative_to(input_root)).replace("/", "\\")
    except ValueError:
        rel_path = file_path.name
    manifest_row = manifest.get(rel_path.lower(), {})
    filename_text = file_path.stem
    title = (manifest_row.get("source_title") or filename_text).strip()
    question_hint = (manifest_row.get("question_hint") or title).strip()
    source_org = (manifest_row.get("source_org") or "").strip()
    category = (manifest_row.get("category") or "通用").strip()
    subcategory = (manifest_row.get("subcategory") or infer_subcategory(file_path)).strip()
    lab_type = (manifest_row.get("lab_type") or infer_lab_type(title + rel_path)).strip()
    risk_level = (manifest_row.get("risk_level") or "3").strip()
    hazard_types = (manifest_row.get("hazard_types") or infer_hazards(title + rel_path)).strip()
    tags = (manifest_row.get("tags") or f"{file_path.suffix.lower().lstrip('.')};{file_path.parent.name}").strip()
    language = (manifest_row.get("language") or "zh-CN").strip()
    reviewer = (manifest_row.get("reviewer") or "").strip()
    return {
        "path": rel_path,
        "source_title": title,
        "source_org": source_org,
        "category": category,
        "subcategory": subcategory,
        "lab_type": lab_type,
        "risk_level": risk_level,
        "hazard_types": hazard_types,
        "tags": tags,
        "language": language,
        "question_hint": question_hint,
        "reviewer": reviewer,
    }


def extract_pdf(path: Path, ocr_mode: str = "auto", special_rule: dict | None = None) -> tuple[str, dict]:
    candidates: list[dict] = []
    candidate_errors: list[dict] = []
    ocr_candidate: dict | None = None
    effective_ocr_mode = "always" if special_rule and special_rule.get("force_ocr") else ocr_mode

    try:
        primary = build_pdf_candidate("pypdf", extract_pdf_pages_pypdf(path), special_rule=special_rule)
        candidates.append(primary)
    except Exception as exc:
        primary = None
        candidate_errors.append({"method": "pypdf", "error": str(exc)})

    if primary is None or should_fallback_from_pypdf(primary):
        for method, extractor in (
            ("pymupdf", extract_pdf_pages_pymupdf),
            ("pdfplumber", extract_pdf_pages_pdfplumber),
        ):
            try:
                candidates.append(build_pdf_candidate(method, extractor(path), special_rule=special_rule))
            except Exception as exc:
                candidate_errors.append({"method": method, "error": str(exc)})

    if not candidates:
        errors = "; ".join(f"{item['method']}: {item['error']}" for item in candidate_errors)
        raise RuntimeError(errors or f"Failed to extract PDF: {path.name}")

    best_candidate = choose_best_pdf_candidate(candidates)

    if should_try_ocr(best_candidate, effective_ocr_mode):
        try:
            ocr_candidate = build_pdf_candidate("ocr", extract_pdf_pages_ocr(path), special_rule=special_rule)
            candidates.append(ocr_candidate)
            best_candidate = choose_best_pdf_candidate(candidates)
        except Exception as exc:
            candidate_errors.append({"method": "ocr", "error": str(exc)})

    final_text = best_candidate["text"]
    ocr_title_patched = False
    if ocr_candidate and best_candidate["method"] != "ocr":
        final_text, ocr_title_patched = patch_pdf_title_from_ocr(best_candidate["text"], ocr_candidate["text"])

    payload = dict(best_candidate["clean_meta"])
    payload["pdf_extractor"] = best_candidate["method"]
    payload["pdf_candidate_summaries"] = [item["summary"] for item in candidates]
    if ocr_title_patched:
        payload["pdf_ocr_title_patched"] = True
    if candidate_errors:
        payload["pdf_candidate_errors"] = candidate_errors

    return final_text, payload


def extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_lines.append(text)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
    return "\n\n".join(parts)


def find_soffice() -> str | None:
    return shutil.which("soffice") or shutil.which("libreoffice")


def convert_legacy_file(path: Path, conversion_dir: Path) -> Path:
    target_suffix = CONVERSION_TARGETS[path.suffix.lower()]
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError(f"Legacy file {path.name} requires LibreOffice/soffice for conversion")
    conversion_dir.mkdir(parents=True, exist_ok=True)
    command = [
        soffice,
        "--headless",
        "--convert-to",
        target_suffix.lstrip("."),
        "--outdir",
        str(conversion_dir),
        str(path),
    ]
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    if completed.returncode != 0:
        raise RuntimeError(completed.stderr.strip() or completed.stdout.strip() or f"Failed to convert {path.name}")
    converted = conversion_dir / f"{path.stem}{target_suffix}"
    if not converted.exists():
        raise RuntimeError(f"Converted file not found for {path.name}")
    return converted


def extract_text_from_file(
    path: Path,
    conversion_dir: Path,
    pdf_ocr_mode: str,
    pdf_special_rule: dict | None = None,
) -> tuple[str, str, Path, dict]:
    effective_path = path
    suffix = path.suffix.lower()
    if suffix in CONVERSION_TARGETS:
        effective_path = convert_legacy_file(path, conversion_dir)
        suffix = effective_path.suffix.lower()
    if suffix == ".pdf":
        text, extraction_meta = extract_pdf(
            effective_path,
            ocr_mode=pdf_ocr_mode,
            special_rule=pdf_special_rule,
        )
        return text, "PDF", effective_path, extraction_meta
    if suffix == ".docx":
        return extract_docx(effective_path), "Word", effective_path, {}
    if suffix == ".pptx":
        return extract_pptx(effective_path), "PPT", effective_path, {}
    raise RuntimeError(f"Unsupported file type: {path.suffix}")


def discover_input_files(input_root: Path, extract_zips: bool, extracted_dir: Path) -> list[Path]:
    files: list[Path] = []
    seen: set[str] = set()

    def add_candidate(path: Path) -> None:
        normalized = str(path.resolve()).lower()
        if normalized in seen:
            return
        seen.add(normalized)
        files.append(path)

    for path in input_root.rglob("*"):
        if not path.is_file():
            continue
        suffix = path.suffix.lower()
        if suffix in SUPPORTED_EXTENSIONS:
            add_candidate(path)
        elif extract_zips and suffix == ZIP_EXTENSION:
            target_dir = extracted_dir / path.stem
            if not target_dir.exists():
                with zipfile.ZipFile(path, "r") as zip_ref:
                    zip_ref.extractall(target_dir)
            for extracted_file in target_dir.rglob("*"):
                if extracted_file.is_file() and extracted_file.suffix.lower() in SUPPORTED_EXTENSIONS:
                    add_candidate(extracted_file)
    return sorted(files)


def build_extract_record(
    file_path: Path,
    metadata: dict,
    text: str,
    source_type: str,
    effective_path: Path,
    extraction_meta: dict,
) -> dict:
    stat = effective_path.stat()
    normalized = normalize_text(text)
    preview = normalized[:240]
    return {
        "source_id": metadata["path"],
        "file_path": str(file_path),
        "effective_path": str(effective_path),
        "source_title": metadata["source_title"],
        "source_org": metadata["source_org"],
        "category": metadata["category"],
        "subcategory": metadata["subcategory"],
        "lab_type": metadata["lab_type"],
        "risk_level": metadata["risk_level"],
        "hazard_types": metadata["hazard_types"],
        "tags": metadata["tags"],
        "language": metadata["language"],
        "question_hint": metadata["question_hint"],
        "reviewer": metadata["reviewer"],
        "source_type": source_type,
        "source_date": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d"),
        "content": normalized,
        "content_preview": preview,
        "content_sha1": sha1_text(normalized) if normalized else "",
        "char_count": len(normalized),
        "extraction_meta": extraction_meta,
    }


def build_kb_rows(documents: list[dict], max_chars: int, overlap: int) -> list[dict]:
    today = datetime.now().strftime("%Y-%m-%d")
    rows: list[dict] = []
    for document in documents:
        chunks = split_into_chunks(document["content"], max_chars=max_chars, overlap=overlap)
        question_hint = document["question_hint"] or document["source_title"]
        question = (
            question_hint
            if any(token in question_hint for token in ["哪些", "什么", "如何", "是否", "能否", "吗", "?", "？"])
            else f"{question_hint}有哪些要点？"
        )
        for idx, chunk in enumerate(chunks, start=1):
            rows.append(
                {
                    "id": f"DOC-{sha1_text(document['file_path'])[:8]}-{idx:03d}",
                    "title": document["source_title"] if idx == 1 else f"{document['source_title']}（片段{idx}）",
                    "category": document["category"],
                    "subcategory": document["subcategory"],
                    "lab_type": document["lab_type"],
                    "risk_level": document["risk_level"],
                    "hazard_types": document["hazard_types"],
                    "scenario": question_hint,
                    "question": question,
                    "answer": chunk,
                    "steps": "",
                    "ppe": "",
                    "forbidden": "",
                    "disposal": "",
                    "first_aid": "",
                    "emergency": "",
                    "legal_notes": "基于本地文档自动抽取，导入知识库前请人工复核。",
                    "references": f"{document['source_title']} | {document['file_path']}",
                    "source_type": document["source_type"],
                    "source_title": document["source_title"],
                    "source_org": document["source_org"],
                    "source_version": "",
                    "source_date": document["source_date"],
                    "source_url": document["file_path"],
                    "last_updated": today,
                    "reviewer": document["reviewer"],
                    "status": "draft",
                    "tags": document["tags"],
                    "language": document["language"],
                }
            )
    return rows


def write_report(path: Path, discovered: int, extracted: int, documents: int, rows: int, skipped: list[dict]) -> None:
    payload = {
        "generated_at": now_iso(),
        "discovered_files": discovered,
        "extracted_successfully": extracted,
        "clean_documents": documents,
        "knowledge_rows": rows,
        "skipped_files": skipped,
    }
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Ingest PDF/Word/PPT files into knowledge-base CSV.")
    parser.add_argument(
        "--input-root",
        default="..\\data",
        help="Root directory that contains documents or extracted archives.",
    )
    parser.add_argument(
        "--manifest",
        default="data_sources\\document_manifest.csv",
        help="Optional CSV that overrides metadata per file. Use the template if you do not have one yet.",
    )
    parser.add_argument(
        "--only-manifest",
        action="store_true",
        help="Only ingest files that are explicitly listed in the manifest.",
    )
    parser.add_argument(
        "--pdf-special-rules",
        default="data_sources\\pdf_special_rules.csv",
        help="Optional CSV for per-file PDF extraction overrides such as force_ocr and body_start_page.",
    )
    parser.add_argument(
        "--output-dir",
        default="artifacts\\document_ingest",
        help="Directory for extracted metadata, cleaned docs, and exported CSV.",
    )
    parser.add_argument(
        "--extract-zips",
        action="store_true",
        help="Extract zip archives found under input-root and scan their contents.",
    )
    parser.add_argument(
        "--max-chars",
        type=int,
        default=1400,
        help="Maximum characters per exported knowledge chunk.",
    )
    parser.add_argument(
        "--overlap",
        type=int,
        default=120,
        help="Character overlap between chunks.",
    )
    parser.add_argument(
        "--limit",
        type=int,
        default=0,
        help="Optional cap on number of discovered files for testing.",
    )
    parser.add_argument(
        "--merge-into",
        default="",
        help="Optional path to append unique rows into an existing knowledge-base CSV.",
    )
    parser.add_argument(
        "--pdf-ocr-mode",
        choices=["auto", "off", "always"],
        default="auto",
        help="PDF fallback mode: auto uses PyMuPDF/pdfplumber/OCR only when needed.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    input_root = Path(args.input_root).resolve()
    output_dir = Path(args.output_dir)
    extracted_dir = output_dir / "extracted_archives"
    conversion_dir = output_dir / "converted_legacy"
    output_dir.mkdir(parents=True, exist_ok=True)

    manifest_path = Path(args.manifest)
    manifest = load_manifest(manifest_path if manifest_path.exists() else None)
    if args.only_manifest and not manifest:
        raise SystemExit("--only-manifest was set, but the manifest file was not found or was empty.")
    pdf_special_rules_path = Path(args.pdf_special_rules)
    pdf_special_rules = load_pdf_special_rules(
        pdf_special_rules_path if pdf_special_rules_path.exists() else None
    )

    discovered_files = discover_input_files(
        input_root=input_root,
        extract_zips=args.extract_zips,
        extracted_dir=extracted_dir,
    )
    if args.only_manifest:
        manifest_paths = set(manifest.keys())
        discovered_files = [
            file_path
            for file_path in discovered_files
            if normalize_rel_path(
                str(file_path.relative_to(input_root)) if file_path.is_relative_to(input_root) else file_path.name
            )
            in manifest_paths
        ]
    if args.limit and args.limit > 0:
        discovered_files = discovered_files[: args.limit]

    extract_results: list[dict] = []
    clean_documents: list[dict] = []
    skipped_files: list[dict] = []

    for file_path in discovered_files:
        metadata = build_metadata(file_path, input_root=input_root, manifest=manifest)
        pdf_special_rule = (
            match_pdf_special_rule(metadata["path"], pdf_special_rules)
            if file_path.suffix.lower() == ".pdf"
            else None
        )
        try:
            text, source_type, effective_path, extraction_meta = extract_text_from_file(
                file_path,
                conversion_dir,
                args.pdf_ocr_mode,
                pdf_special_rule=pdf_special_rule,
            )
            record = build_extract_record(
                file_path,
                metadata,
                text,
                source_type,
                effective_path,
                extraction_meta,
            )
            extract_results.append(record)
            if len(record["content"]) >= 80:
                clean_documents.append(record)
            else:
                skipped_files.append({"file_path": str(file_path), "reason": "content_too_short"})
        except Exception as exc:
            skipped_files.append({"file_path": str(file_path), "reason": str(exc)})

    write_jsonl(output_dir / "extract_results.jsonl", extract_results)
    write_jsonl(output_dir / "clean_documents.jsonl", clean_documents)

    kb_rows = build_kb_rows(clean_documents, max_chars=args.max_chars, overlap=args.overlap)
    write_csv(output_dir / "knowledge_base_documents.csv", kb_rows)

    if args.merge_into:
        merged_count = merge_into_csv(Path(args.merge_into), kb_rows)
        print(f"Merged rows into {Path(args.merge_into).resolve()}: {merged_count}")

    write_report(
        output_dir / "run_report.json",
        discovered=len(discovered_files),
        extracted=len(extract_results),
        documents=len(clean_documents),
        rows=len(kb_rows),
        skipped=skipped_files,
    )

    print(f"Discovered files: {len(discovered_files)}")
    print(f"Extracted successfully: {len(extract_results)}")
    print(f"Clean documents: {len(clean_documents)}")
    print(f"Knowledge rows: {len(kb_rows)}")
    print(f"Output dir: {output_dir.resolve()}")


if __name__ == "__main__":
    main()
